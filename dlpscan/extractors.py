"""Text extraction from binary document formats.

Extracts plain text from Office documents, PDFs, and email files so that
the dlpscan scanner can process them. All extraction libraries are optional
dependencies — a clear error message is raised if the required library is
not installed.

Usage::

    from dlpscan.extractors import extract_text

    result = extract_text('report.pdf')
    print(result.text[:200])
    print(result.metadata)

Install extraction dependencies::

    pip install dlpscan[pdf]          # PDF support
    pip install dlpscan[office]       # DOCX, XLSX, PPTX support
    pip install dlpscan[email]        # MSG support (EML uses stdlib)
    pip install dlpscan[all-formats]  # Everything
"""

import email
import email.policy
import logging
import os
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, List, Optional

from .exceptions import ExtractionError

logger = logging.getLogger(__name__)

# Maximum file size for extraction (100 MB default).
MAX_EXTRACT_SIZE = 100 * 1024 * 1024


@dataclass
class ExtractionResult:
    """Result of extracting text from a document.

    Attributes:
        text: Extracted plain text content.
        metadata: Format-specific metadata (author, page count, etc.).
        format: Detected format identifier (e.g., 'pdf', 'docx').
        warnings: Non-fatal issues encountered during extraction.
    """
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    format: str = 'unknown'
    warnings: List[str] = field(default_factory=list)


# -- Extractor registry --
# Maps lowercase file extension (with dot) to extractor function.
_EXTRACTORS: Dict[str, Callable[[str], ExtractionResult]] = {}


def register_extractor(extension: str, func: Callable[[str], ExtractionResult]) -> None:
    """Register a custom text extractor for a file extension.

    Args:
        extension: File extension with dot (e.g., '.pdf').
        func: Callable that takes a file path and returns ExtractionResult.

    Example::

        def extract_rtf(path):
            text = my_rtf_parser(path)
            return ExtractionResult(text=text, format='rtf')

        register_extractor('.rtf', extract_rtf)
    """
    if not extension.startswith('.'):
        raise ValueError("Extension must start with '.' (e.g., '.pdf')")
    if not callable(func):
        raise TypeError("func must be callable.")
    _EXTRACTORS[extension.lower()] = func


def get_extractor(file_path: str) -> Optional[Callable[[str], ExtractionResult]]:
    """Return the registered extractor for a file, or None."""
    _, ext = os.path.splitext(file_path)
    return _EXTRACTORS.get(ext.lower())


def extract_text(file_path: str, max_size: int = MAX_EXTRACT_SIZE) -> ExtractionResult:
    """Extract text from any supported file format.

    Routes to the appropriate extractor based on file extension.
    Falls back to plain text reading for unrecognized extensions.

    Args:
        file_path: Path to the file.
        max_size: Maximum file size in bytes (default 100 MB).

    Returns:
        ExtractionResult with extracted text and metadata.

    Raises:
        FileNotFoundError: If the file does not exist.
        ExtractionError: If extraction fails.
        ValueError: If the file exceeds max_size.
    """
    try:
        file_size = os.path.getsize(file_path)
    except OSError:
        raise FileNotFoundError(f"File not found: {file_path}")

    if file_size > max_size:
        raise ValueError(
            f"File exceeds maximum size of {max_size:,} bytes "
            f"({file_size:,} bytes)."
        )

    if file_size == 0:
        return ExtractionResult(text='', format='empty', metadata={'size': 0})

    extractor = get_extractor(file_path)
    if extractor is not None:
        return extractor(file_path)

    # Fallback: treat as plain text.
    return _extract_plain_text(file_path)


def supported_extensions() -> List[str]:
    """Return a sorted list of all registered file extensions."""
    return sorted(_EXTRACTORS.keys())


# -- Built-in extractors --

def _extract_plain_text(file_path: str) -> ExtractionResult:
    """Read a plain text file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
    except OSError as exc:
        raise ExtractionError(f"Failed to read text file: {exc}")
    return ExtractionResult(
        text=text,
        format='text',
        metadata={'size': len(text)},
    )


def _extract_pdf(file_path: str) -> ExtractionResult:
    """Extract text from a PDF file using pdfplumber with OCR fallback.

    Pages that yield no text via pdfplumber are automatically processed
    with OCR if pytesseract is available.
    """
    try:
        import pdfplumber
    except ImportError:
        raise ExtractionError(
            "PDF extraction requires 'pdfplumber'. "
            "Install with: pip install dlpscan[pdf]"
        )

    warnings = []
    pages_text = []
    empty_pages = []
    metadata = {}

    try:
        with pdfplumber.open(file_path) as pdf:
            metadata = {
                'page_count': len(pdf.pages),
                'pdf_metadata': pdf.metadata or {},
            }
            for i, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        pages_text.append((i, page_text))
                    else:
                        empty_pages.append(i)
                except Exception as exc:
                    warnings.append(f"Page {i + 1}: extraction failed ({exc})")
                    empty_pages.append(i)
    except Exception as exc:
        raise ExtractionError(f"Failed to open PDF: {exc}")

    # OCR fallback for pages that yielded no text.
    if empty_pages:
        try:
            from .ocr import ocr_available
            if ocr_available():
                from pdf2image import convert_from_path

                from .ocr import DEFAULT_CONFIG, DEFAULT_LANG, _ensure_pytesseract
                pytesseract = _ensure_pytesseract()
                for page_idx in empty_pages:
                    try:
                        images = convert_from_path(
                            file_path, dpi=300,
                            first_page=page_idx + 1,
                            last_page=page_idx + 1,
                            thread_count=1,
                        )
                        if images:
                            page_text = pytesseract.image_to_string(
                                images[0], lang=DEFAULT_LANG, config=DEFAULT_CONFIG
                            ).strip()
                            if page_text:
                                pages_text.append((page_idx, page_text))
                                metadata.setdefault('ocr_pages', []).append(page_idx + 1)
                    except Exception as exc:
                        warnings.append(f"Page {page_idx + 1}: OCR fallback failed ({exc})")
            else:
                if empty_pages:
                    warnings.append(
                        f"{len(empty_pages)} page(s) had no extractable text. "
                        "Install dlpscan[ocr] for OCR fallback."
                    )
        except ImportError:
            warnings.append(
                f"{len(empty_pages)} page(s) had no extractable text. "
                "Install dlpscan[ocr] for OCR fallback."
            )

    # Sort pages by index and join.
    pages_text.sort(key=lambda x: x[0])
    combined = '\n\n'.join(text for _, text in pages_text)

    return ExtractionResult(
        text=combined,
        format='pdf',
        metadata=metadata,
        warnings=warnings,
    )


def _extract_docx(file_path: str) -> ExtractionResult:
    """Extract text from a DOCX file using python-docx."""
    try:
        import docx
    except ImportError:
        raise ExtractionError(
            "DOCX extraction requires 'python-docx'. "
            "Install with: pip install dlpscan[office]"
        )

    try:
        doc = docx.Document(file_path)
    except Exception as exc:
        raise ExtractionError(f"Failed to open DOCX: {exc}")

    parts = []

    # Extract paragraphs.
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Extract table cells.
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                parts.append('\t'.join(row_text))

    # Metadata from core properties.
    metadata = {}
    try:
        props = doc.core_properties
        metadata = {
            'author': props.author or '',
            'title': props.title or '',
            'created': str(props.created) if props.created else '',
            'modified': str(props.modified) if props.modified else '',
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables),
        }
    except Exception:
        pass

    return ExtractionResult(
        text='\n'.join(parts),
        format='docx',
        metadata=metadata,
    )


def _extract_xlsx(file_path: str) -> ExtractionResult:
    """Extract text from an XLSX file using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        raise ExtractionError(
            "XLSX extraction requires 'openpyxl'. "
            "Install with: pip install dlpscan[office]"
        )

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(f"Failed to open XLSX: {exc}")

    parts = []
    warnings = []
    sheet_names = wb.sheetnames

    try:
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            try:
                for row in ws.iter_rows(values_only=True):
                    cell_texts = [str(cell) for cell in row if cell is not None]
                    if cell_texts:
                        parts.append('\t'.join(cell_texts))
            except Exception as exc:
                warnings.append(f"Sheet '{sheet_name}': read error ({exc})")
    finally:
        wb.close()

    return ExtractionResult(
        text='\n'.join(parts),
        format='xlsx',
        metadata={'sheet_count': len(sheet_names), 'sheets': sheet_names},
        warnings=warnings,
    )


def _extract_pptx(file_path: str) -> ExtractionResult:
    """Extract text from a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ExtractionError(
            "PPTX extraction requires 'python-pptx'. "
            "Install with: pip install dlpscan[office]"
        )

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise ExtractionError(f"Failed to open PPTX: {exc}")

    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_parts.append('\t'.join(row_text))
        if slide_parts:
            parts.append('\n'.join(slide_parts))

    return ExtractionResult(
        text='\n\n'.join(parts),
        format='pptx',
        metadata={'slide_count': len(prs.slides)},
    )


def _extract_eml(file_path: str) -> ExtractionResult:
    """Extract text from an EML file using the stdlib email module."""
    try:
        with open(file_path, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=email.policy.default)
    except Exception as exc:
        raise ExtractionError(f"Failed to parse EML: {exc}")

    parts = []
    warnings = []

    # Headers.
    headers = {}
    for header in ('From', 'To', 'Cc', 'Subject', 'Date'):
        value = msg.get(header, '')
        if value:
            headers[header.lower()] = str(value)
            parts.append(f"{header}: {value}")

    # Body parts.
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == 'text/plain':
                try:
                    body = part.get_content()
                    if isinstance(body, str) and body.strip():
                        parts.append(body)
                except Exception as exc:
                    warnings.append(f"Part decode error: {exc}")
            elif content_type == 'text/html':
                try:
                    body = part.get_content()
                    if isinstance(body, str) and body.strip():
                        # Strip HTML tags for basic text extraction.
                        import re
                        clean = re.sub(r'<[^>]+>', ' ', body)
                        clean = re.sub(r'\s+', ' ', clean).strip()
                        if clean:
                            parts.append(clean)
                except Exception as exc:
                    warnings.append(f"HTML part decode error: {exc}")
    else:
        try:
            body = msg.get_content()
            if isinstance(body, str) and body.strip():
                parts.append(body)
        except Exception as exc:
            warnings.append(f"Body decode error: {exc}")

    return ExtractionResult(
        text='\n\n'.join(parts),
        format='eml',
        metadata={
            'headers': headers,
            'is_multipart': msg.is_multipart(),
        },
        warnings=warnings,
    )


def _extract_msg(file_path: str) -> ExtractionResult:
    """Extract text from an Outlook MSG file using extract-msg."""
    try:
        import extract_msg
    except ImportError:
        raise ExtractionError(
            "MSG extraction requires 'extract-msg'. "
            "Install with: pip install dlpscan[email]"
        )

    try:
        msg = extract_msg.openMsg(file_path)
    except Exception as exc:
        raise ExtractionError(f"Failed to open MSG: {exc}")

    parts = []
    try:
        headers = {}
        for attr, label in [('sender', 'From'), ('to', 'To'), ('cc', 'Cc'),
                            ('subject', 'Subject'), ('date', 'Date')]:
            value = getattr(msg, attr, None)
            if value:
                headers[label.lower()] = str(value)
                parts.append(f"{label}: {value}")

        body = msg.body
        if body and body.strip():
            parts.append(body)
    finally:
        msg.close()

    return ExtractionResult(
        text='\n\n'.join(parts),
        format='msg',
        metadata={'headers': headers},
    )


def _extract_image_ocr(file_path: str) -> ExtractionResult:
    """Extract text from an image file using OCR."""
    try:
        from .ocr import ocr_available, ocr_image
    except ImportError:
        raise ExtractionError(
            "Image OCR requires 'pytesseract' and 'Pillow'. "
            "Install with: pip install dlpscan[ocr]"
        )

    if not ocr_available():
        raise ExtractionError(
            "OCR is not available. Ensure Tesseract is installed: "
            "apt install tesseract-ocr (Linux) or brew install tesseract (macOS). "
            "Then: pip install dlpscan[ocr]"
        )

    try:
        result = ocr_image(file_path)
    except Exception as exc:
        raise ExtractionError(f"Image OCR failed: {exc}")

    return ExtractionResult(
        text=result.text,
        format='image_ocr',
        metadata={
            'ocr_confidence': result.confidence,
            'ocr_language': result.language,
            **result.metadata,
        },
        warnings=result.warnings,
    )


def _extract_legacy_office(file_path: str) -> ExtractionResult:
    """Placeholder for legacy Office formats (.doc, .xls, .ppt).

    These binary formats require external tools (e.g., LibreOffice, antiword,
    catdoc) for text extraction. Raises a clear error with guidance.
    """
    _, ext = os.path.splitext(file_path)
    raise ExtractionError(
        f"Legacy Office format '{ext}' requires external tools for text extraction. "
        "Options: (1) Convert to modern format (.docx/.xlsx/.pptx) using LibreOffice, "
        "(2) Use 'textract' package, or "
        "(3) Install LibreOffice and use: soffice --convert-to <format> <file>"
    )


# -- Register built-in extractors --

_EXTRACTORS['.pdf'] = _extract_pdf
_EXTRACTORS['.docx'] = _extract_docx
_EXTRACTORS['.xlsx'] = _extract_xlsx
_EXTRACTORS['.pptx'] = _extract_pptx
_EXTRACTORS['.eml'] = _extract_eml
_EXTRACTORS['.msg'] = _extract_msg
_EXTRACTORS['.doc'] = _extract_legacy_office
_EXTRACTORS['.xls'] = _extract_legacy_office
_EXTRACTORS['.ppt'] = _extract_legacy_office

# Image formats — OCR extraction.
for _ext in ('.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.webp'):
    _EXTRACTORS[_ext] = _extract_image_ocr

# Common plain text extensions get explicit registration so they bypass
# the binary-file heuristic in scan_directory.
for _ext in ('.txt', '.csv', '.log', '.json', '.xml', '.html', '.htm',
             '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.md',
             '.rst', '.py', '.js', '.ts', '.java', '.go', '.rs', '.rb',
             '.php', '.sh', '.bat', '.ps1', '.sql', '.env'):
    _EXTRACTORS[_ext] = _extract_plain_text
